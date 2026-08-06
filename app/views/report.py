import streamlit as st
import io
import json

def render(opp_df, raw_df):
    st.title("Executive Report Generation")
    
    st.markdown("Generate comprehensive board-ready reports containing the AI Strategy Copilot's findings.")
    
    col1, col2 = st.columns([1, 1])
    
    with col1:
        st.markdown("### Formal Exports")
        report_type = st.radio("Select Format", ["Excel Data Pack", "PDF Summary", "PowerPoint Deck"])
        if st.button("Generate Export"):
            if report_type == "Excel Data Pack":
                csv = opp_df.to_csv(index=False).encode('utf-8')
                st.download_button("Download Excel", data=csv, file_name='CardioScope_V2_DataPack.csv', mime='text/csv')
                
            elif report_type == "PDF Summary":
                try:
                    from fpdf import FPDF
                except ImportError:
                    st.error("Please install fpdf2")
                    return
                pdf = FPDF()
                pdf.add_page()
                pdf.set_font("Helvetica", size=18, style='B')
                pdf.cell(200, 15, txt="CardioScope AI: Executive Strategy Report", ln=1, align='C')
                pdf.set_font("Helvetica", size=14, style='B')
                pdf.cell(200, 10, txt="Top 10 Strategic Priorities", ln=1)
                pdf.set_font("Helvetica", size=10)
                for idx, row in opp_df.head(10).iterrows():
                    pdf.set_font("Helvetica", size=11, style='B')
                    pdf.cell(0, 8, txt=f"{idx+1}. {row['MOLECULE_DESC']} (Strategy: {row.get('Investment Matrix', row['Strategy'])})", ln=1)
                    pdf.set_font("Helvetica", size=10)
                    pdf.cell(0, 6, txt=f"   Opp Score: {row['Opportunity Score']} | Confidence: {row['Confidence Score']}%", ln=1)
                    why_text = "High opportunity."
                    try:
                        exp = json.loads(row.get('Explanation', '{}'))
                        if 'Why' in exp and len(exp['Why']) > 0: why_text = exp['Why'][0]
                    except: pass
                    pdf.multi_cell(0, 6, txt=f"   AI Rationale: {why_text}")
                pdf_bytes = pdf.output(dest='S')
                st.download_button("Download PDF", data=bytes(pdf_bytes), file_name='CardioScope_V2_Report.pdf', mime='application/pdf')
                
            elif report_type == "PowerPoint Deck":
                try:
                    from pptx import Presentation
                    from pptx.util import Inches, Pt
                except ImportError:
                    st.error("Please install python-pptx")
                    return
                prs = Presentation()
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = "CardioScope AI V2"
                subtitle.text = "Executive Strategy Recommendations"
                
                bullet_slide_layout = prs.slide_layouts[1]
                for idx, row in opp_df.head(5).iterrows():
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = slide.shapes
                    title_shape = shapes.title
                    body_shape = shapes.placeholders[1]
                    title_shape.text = f"{idx+1}. {row['MOLECULE_DESC']}"
                    tf = body_shape.text_frame
                    tf.text = f"Strategy: {row.get('Investment Matrix', row['Strategy'])}"
                    p = tf.add_paragraph()
                    p.text = f"Opportunity Score: {row['Opportunity Score']}"
                    p = tf.add_paragraph()
                    p.text = f"Risk Score: {row['Risk Score']}"
                
                ppt_stream = io.BytesIO()
                prs.save(ppt_stream)
                st.download_button("Download PPTX", data=ppt_stream.getvalue(), file_name='CardioScope_V2_Deck.pptx', mime='application/vnd.openxmlformats-officedocument.presentationml.presentation')

    with col2:
        st.markdown("### Executive Brief Generator")
        brief_type = st.selectbox("Select Brief Type", ["CEO Brief", "Strategy Brief", "Investment Memo"])
        if st.button("Generate Brief"):
            top = opp_df.iloc[0]
            if brief_type == "CEO Brief":
                text = f"MEMORANDUM FOR THE CEO\n\nSUBJECT: Top Cardiac Opportunity - {top['MOLECULE_DESC']}\n\nOur AI Copilot strongly recommends a '{top.get('Investment Matrix', top['Strategy'])}' approach for {top['MOLECULE_DESC']}. It scores {top['Opportunity Score']}/100 with a Confidence of {top['Confidence Score']}%.\n\nThe strategic fit is {top.get('Strategic Fit Score', 0)}/100, driven by a strong Market Score ({top['Market Score']}/10) and favorable Future Outlook ({top.get('Future Score', 0)}/10). Immediate resource allocation is advised."
            elif brief_type == "Strategy Brief":
                text = f"STRATEGY BRIEF: {top['MOLECULE_DESC']}\n\nPrimary Directive: {top.get('Investment Matrix', top['Strategy'])}\n\nKey Metrics:\n- Opportunity Index: {top['Opportunity Score']}\n- HHI Market Concentration: {top.get('HHI', 0):.0f}\n- Risk Score: {top['Risk Score']}/10\n\nAction Plan: Maximize Cipla's Right-to-Win ({top['Right to Win Score']}/10) while monitoring patent threats."
            else:
                text = f"INVESTMENT MEMO\n\nAsset: {top['MOLECULE_DESC']}\nRecommendation: {top.get('Investment Matrix', top['Strategy'])}\n\nWe calculate a Business Impact Score of {top.get('Business Impact Score', 0)}/100. The Entry Barrier Score is {top.get('Entry Barrier Score', 0)}/10, suggesting defensibility. Deploy capital aggressively to capture Growth ({top['Growth Score']}/10)."
            
            st.text_area("Generated Output (Editable)", text, height=200)
            st.download_button(f"Download {brief_type}", data=text.encode('utf-8'), file_name=f"{brief_type.replace(' ', '_')}.txt", mime='text/plain')
